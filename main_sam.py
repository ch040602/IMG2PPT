# -*- coding: utf-8 -*-
"""
SlimSAM + AutomaticMaskGenerator 기반 PPT 도형 변환기
(https://github.com/czg1225/SlimSAM)

로드 방식
---------
```python
model_path = "checkpoints/SlimSAM-77.pth"
SlimSAM_model = torch.load(model_path, map_location=device)
SlimSAM_model.image_encoder = SlimSAM_model.image_encoder.module  # DDP 래핑 해제
SlimSAM_model.to(device)
SlimSAM_model.eval()
```
→ 그 후 `SamAutomaticMaskGenerator` 사용
"""

import os
import types
import cv2
import numpy as np
import torch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from tkinter import filedialog, Tk, Button, Checkbutton, IntVar, Label
from segment_anything import SamAutomaticMaskGenerator
from segment_anything import SamPredictor
# ---------------- GUI ----------------
root = Tk()
root.title("SlimSAM AutoMask → PPT")
root.geometry("440x240")

image_path: str | None = None
outline_var = IntVar()
Checkbutton(root, text="도형 테두리 포함", variable=outline_var).pack()
Label(root, text="1) Load Image → 2) Convert").pack(pady=5)

# ---------------- SlimSAM 로드 ----------------
DEVICE = torch.device("cuda" if torch.cuda.is_available() else "cpu")
MODEL_PATH = "./SlimSAM/checkpoints/SlimSAM-77.pth"  # 유저 지정 경로 그대로 사용

print("[SlimSAM] loading", MODEL_PATH)
# 모델 로드
device = torch.device("cuda")
SlimSAM_model = torch.load(MODEL_PATH)

SlimSAM_model.image_encoder = SlimSAM_model.image_encoder.module
SlimSAM_model.to(device)
SlimSAM_model.eval()

def forward(self, x):
    x = self.patch_embed(x)
    if self.pos_embed is not None:
        x = x + self.pos_embed
    for blk in self.blocks:
        x, qkv_emb, mid_emb, x_emb = blk(x)
    x = self.neck(x.permute(0, 3, 1, 2))
    return x

import types
SlimSAM_model.image_encoder.forward = types.MethodType(forward, SlimSAM_model.image_encoder)

predictor = SamPredictor(SlimSAM_model)

# SAM 설정: 세부 조각까지 최대한 유지하도록 파라미터 완화
mask_generator = SamAutomaticMaskGenerator(
    model=SlimSAM_model,
    points_per_side=64,
    pred_iou_thresh=0.6,
    stability_score_thresh=0.6,
    crop_n_layers=2,
    crop_n_points_downscale_factor=1,
    min_mask_region_area=5,
)

# ---------------- Helper ----------------

def load_image():
    global image_path
    image_path = filedialog.askopenfilename(filetypes=[("Image", "*.png;*.jpg;*.jpeg")])
    if image_path:
        root.title("선택됨: " + os.path.basename(image_path))


def convert_to_ppt():
    global image_path
    if not image_path:
        print("이미지 먼저 선택")
        return

    bgr = cv2.imread(image_path)
    if bgr is None:
        print("로드 실패")
        return
    rgb = cv2.cvtColor(bgr, cv2.COLOR_BGR2RGB)

    print("[SAM] 마스크 생성 중 …")
    masks = mask_generator.generate(rgb)
    print(f"[SAM] 마스크 개수: {len(masks)}")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    h, w = rgb.shape[:2]

    for m in masks:
        mask = m["segmentation"].astype(np.uint8) * 255
        if cv2.countNonZero(mask) < 5:
            continue
        cnts, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_NONE)
        if not cnts:
            continue
        cnt = max(cnts, key=cv2.contourArea)
        if cv2.contourArea(cnt) < 5:
            continue
        pts = cnt.squeeze()
        if pts.ndim != 2 or len(pts) < 3:
            continue

        mean_col = cv2.mean(rgb, mask=mask)[:3]
        mean_col = tuple(map(int, mean_col))

        builder = slide.shapes.build_freeform()
        x0, y0 = pts[0]
        builder.move_to(Inches(x0 / w * sw), Inches(y0 / h * sh))
        builder.add_line_segments([(Inches(x / w * sw), Inches(y / h * sh)) for x, y in pts[1:]])
        shp = builder.convert_to_shape()
        shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(*mean_col)

        if outline_var.get():
            shp.line.fill.solid(); shp.line.fill.fore_color.rgb = RGBColor(0, 0, 0); shp.line.width = Pt(0.35)
        else:
            shp.line.fill.background()

    save = os.path.splitext(os.path.basename(image_path))[0] + "_sam.pptx"
    prs.save(save)
    print("저장 완료:", save)

# ---------------- Buttons ----------------
Button(root, text="Load Image", command=load_image).pack(pady=6)
Button(root, text="Convert (AutoMask)", command=convert_to_ppt).pack(pady=10)

root.mainloop()