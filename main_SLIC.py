# -*- coding: utf-8 -*-
import os, cv2, numpy as np
from skimage.segmentation import slic
from skimage.util import img_as_float
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from tkinter import filedialog, Tk, Button, Checkbutton, IntVar, Label, Scale, HORIZONTAL

root = Tk()
root.title("SLIC 기반 도형 PPT 변환기 (부드러운 설정)")
root.geometry("420x250")

image_path: str | None = None
outline_var = IntVar()
seg_scale = Scale(root, from_=50, to=400, orient=HORIZONTAL, label="슈퍼픽셀 수 (K)", resolution=25)
seg_scale.set(120); seg_scale.pack(pady=5)
Checkbutton(root, text="도형 테두리 포함", variable=outline_var).pack()
Label(root, text="1) Load → 2) Convert").pack(pady=5)

def load_image():
    global image_path
    image_path = filedialog.askopenfilename(filetypes=[("Image", "*.png;*.jpg;*.jpeg")])
    if image_path:
        root.title("선택됨: " + os.path.basename(image_path))

def convert_to_ppt():
    global image_path
    if not image_path:
        print("이미지 선택 필요")
        return
    img = cv2.imread(image_path)
    if img is None: return
    rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)

    k = seg_scale.get()
    segments = slic(
        img_as_float(rgb),
        n_segments=k,
        compactness=15,        # 공간 중심으로 부드럽게
        sigma=1,               # 경계 블러 적용
        enforce_connectivity=True,
        start_label=0
    )

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    h, w = rgb.shape[:2]
    sw, sh = prs.slide_width.inches, prs.slide_height.inches

    for val in np.unique(segments):
        mask = (segments == val).astype(np.uint8) * 255
        if cv2.countNonZero(mask) < 200:
            continue
        cnts, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        if not cnts: continue
        cnt = max(cnts, key=cv2.contourArea)
        if cv2.contourArea(cnt) < 150: continue
        pts = cnt.squeeze()
        if pts.ndim != 2 or len(pts) < 3: continue

        mean_col = tuple(map(int, cv2.mean(rgb, mask=mask)[:3]))
        builder = slide.shapes.build_freeform()
        x0, y0 = pts[0]
        builder.move_to(Inches(x0 / w * sw), Inches(y0 / h * sh))
        builder.add_line_segments([(Inches(x / w * sw), Inches(y / h * sh)) for x, y in pts[1:]])
        shp = builder.convert_to_shape()
        shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(*mean_col)

        if outline_var.get():
            shp.line.fill.solid(); shp.line.fill.fore_color.rgb = RGBColor(0, 0, 0); shp.line.width = Pt(0.35)
        else:
            shp.line.fill.background()

    save_name = os.path.splitext(os.path.basename(image_path))[0] + "_slic.pptx"
    prs.save(save_name)
    print("저장 완료:", save_name)

Button(root, text="Load Image", command=load_image).pack(pady=5)
Button(root, text="Convert (SLIC 부드럽게)", command=convert_to_ppt).pack(pady=10)
root.mainloop()
